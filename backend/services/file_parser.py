from __future__ import annotations

import csv
import json
import mimetypes
import re
from dataclasses import dataclass
from io import StringIO
from pathlib import Path
from typing import Any

import aiofiles
from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from pypdf import PdfReader

from config.settings import settings
from services.ocr import ocr_service

try:  # pragma: no cover - optional dependency
    import fitz
except Exception:  # pragma: no cover - optional dependency
    fitz = None

try:  # pragma: no cover - optional dependency
    import pdfplumber
except Exception:  # pragma: no cover - optional dependency
    pdfplumber = None

try:  # pragma: no cover - optional dependency
    import pandas as pd
except Exception:  # pragma: no cover - optional dependency
    pd = None

try:  # pragma: no cover - optional dependency
    from openpyxl import load_workbook
except Exception:  # pragma: no cover - optional dependency
    load_workbook = None

try:  # pragma: no cover - optional dependency
    from pptx import Presentation
except Exception:  # pragma: no cover - optional dependency
    Presentation = None


CODE_LANGUAGE_BY_EXTENSION = {
    ".py": "python",
    ".js": "javascript",
    ".jsx": "javascript",
    ".ts": "typescript",
    ".tsx": "typescript",
    ".java": "java",
    ".c": "c",
    ".cpp": "cpp",
    ".cs": "csharp",
    ".go": "go",
    ".rs": "rust",
    ".php": "php",
    ".sql": "sql",
    ".html": "html",
    ".htm": "html",
    ".css": "css",
    ".xml": "xml",
    ".json": "json",
    ".yml": "yaml",
    ".yaml": "yaml",
}


@dataclass
class ParsedSection:
    text: str
    page_number: int | None = None
    sheet_name: str | None = None
    section_title: str | None = None
    language: str | None = None


@dataclass
class ParsedFile:
    text: str
    preview_text: str
    metadata: dict[str, Any]
    sections: list[ParsedSection]


class FileParserService:
    def __init__(self) -> None:
        self.allowed_extensions = set(settings.file_allowed_extensions_list)

    def validate_upload(self, *, filename: str, mime_type: str, size: int) -> tuple[str, str]:
        extension = Path(filename or "").suffix.lower()
        if extension not in self.allowed_extensions:
            raise ValueError(
                "Unsupported file type. Allowed types include PDF, Office docs, text, spreadsheets, images, and code files."
            )
        max_bytes = max(1, int(settings.FILE_MAX_SIZE_MB)) * 1024 * 1024
        if int(size or 0) > max_bytes:
            raise ValueError(f"File exceeds the {settings.FILE_MAX_SIZE_MB}MB limit.")
        detected_mime = mime_type or mimetypes.guess_type(filename or "")[0] or "application/octet-stream"
        return extension, detected_mime

    async def parse(self, path: str, *, original_name: str, mime_type: str = "") -> ParsedFile:
        extension = Path(original_name or path).suffix.lower()
        parser = {
            ".pdf": self._parse_pdf,
            ".docx": self._parse_docx,
            ".txt": self._parse_text,
            ".md": self._parse_text,
            ".csv": self._parse_csv,
            ".xlsx": self._parse_spreadsheet,
            ".xlsm": self._parse_spreadsheet,
            ".xls": self._parse_spreadsheet,
            ".pptx": self._parse_pptx,
            ".png": self._parse_image,
            ".jpg": self._parse_image,
            ".jpeg": self._parse_image,
            ".webp": self._parse_image,
            ".gif": self._parse_image,
            ".bmp": self._parse_image,
        }.get(extension)

        if parser is not None:
            parsed = await parser(path, original_name=original_name, mime_type=mime_type)
        elif extension in CODE_LANGUAGE_BY_EXTENSION:
            parsed = await self._parse_code(path, original_name=original_name, mime_type=mime_type)
        else:
            parsed = await self._parse_text(path, original_name=original_name, mime_type=mime_type)

        text = self._clean_text(parsed.text)
        sections = [
            ParsedSection(
                text=self._clean_text(section.text),
                page_number=section.page_number,
                sheet_name=section.sheet_name,
                section_title=section.section_title,
                language=section.language,
            )
            for section in parsed.sections
            if self._clean_text(section.text)
        ]
        if not sections and text:
            sections = [ParsedSection(text=text)]

        preview_limit = max(200, int(settings.FILE_PREVIEW_CHAR_LIMIT))
        return ParsedFile(
            text=text,
            preview_text=text[:preview_limit],
            metadata=parsed.metadata,
            sections=sections,
        )

    async def _parse_pdf(self, path: str, *, original_name: str, mime_type: str) -> ParsedFile:
        sections: list[ParsedSection] = []
        metadata: dict[str, Any] = {"parser": "pypdf"}

        if fitz is not None:
            document = fitz.open(path)
            try:
                metadata["parser"] = "pymupdf"
                for page_index, page in enumerate(document, start=1):
                    page_text = str(page.get_text("text") or "").strip()
                    if not page_text and ocr_service.available():
                        pixmap = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                        page_text = ocr_service.extract_text_from_bytes(pixmap.tobytes("png"))
                    if page_text:
                        sections.append(ParsedSection(text=page_text, page_number=page_index))
                metadata["page_count"] = len(document)
            finally:
                document.close()
        elif pdfplumber is not None:
            with pdfplumber.open(path) as pdf:
                metadata["parser"] = "pdfplumber"
                metadata["page_count"] = len(pdf.pages)
                for page_index, page in enumerate(pdf.pages, start=1):
                    page_text = str(page.extract_text() or "").strip()
                    if page_text:
                        sections.append(ParsedSection(text=page_text, page_number=page_index))
        else:
            reader = PdfReader(path)
            metadata["page_count"] = len(reader.pages)
            for page_index, page in enumerate(reader.pages, start=1):
                page_text = str(page.extract_text() or "").strip()
                if page_text:
                    sections.append(ParsedSection(text=page_text, page_number=page_index))

        combined = "\n\n".join(section.text for section in sections)
        return ParsedFile(text=combined, preview_text=combined[:400], metadata=metadata, sections=sections)

    async def _parse_docx(self, path: str, *, original_name: str, mime_type: str) -> ParsedFile:
        document = DocxDocument(path)
        sections: list[ParsedSection] = []
        parts: list[str] = []

        for paragraph in document.paragraphs:
            text = str(paragraph.text or "").strip()
            if not text:
                continue
            sections.append(ParsedSection(text=text, section_title="Paragraph"))
            parts.append(text)

        for table_index, table in enumerate(document.tables, start=1):
            rows = []
            for row in table.rows:
                values = [str(cell.text or "").strip() for cell in row.cells]
                if any(values):
                    rows.append(" | ".join(value for value in values if value))
            if rows:
                table_text = "\n".join(rows)
                sections.append(ParsedSection(text=table_text, section_title=f"Table {table_index}"))
                parts.append(f"Table {table_index}\n{table_text}")

        combined = "\n\n".join(parts)
        return ParsedFile(
            text=combined,
            preview_text=combined[:400],
            metadata={"parser": "python-docx"},
            sections=sections,
        )

    async def _parse_text(self, path: str, *, original_name: str, mime_type: str) -> ParsedFile:
        async with aiofiles.open(path, "r", encoding="utf-8", errors="ignore") as file_handle:
            raw_text = await file_handle.read()

        extension = Path(original_name or path).suffix.lower()
        metadata = {"parser": "native", "extension": extension}

        if extension in {".json"}:
            try:
                raw_text = json.dumps(json.loads(raw_text), indent=2)
            except json.JSONDecodeError:
                pass
        elif extension in {".html", ".htm"}:
            raw_text = BeautifulSoup(raw_text, "html.parser").get_text("\n")

        cleaned = self._clean_text(raw_text)
        return ParsedFile(
            text=cleaned,
            preview_text=cleaned[:400],
            metadata=metadata,
            sections=[ParsedSection(text=cleaned)],
        )

    async def _parse_csv(self, path: str, *, original_name: str, mime_type: str) -> ParsedFile:
        if pd is not None:
            data_frame = pd.read_csv(path)
            csv_text = data_frame.to_csv(index=False)
        else:
            async with aiofiles.open(path, "r", encoding="utf-8", errors="ignore") as file_handle:
                raw_content = await file_handle.read()
            reader = csv.reader(StringIO(raw_content))
            rows = [" | ".join(cell.strip() for cell in row if str(cell).strip()) for row in reader]
            csv_text = "\n".join(row for row in rows if row)

        cleaned = self._clean_text(csv_text)
        return ParsedFile(
            text=cleaned,
            preview_text=cleaned[:400],
            metadata={"parser": "pandas" if pd is not None else "csv"},
            sections=[ParsedSection(text=cleaned, sheet_name="CSV")],
        )

    async def _parse_spreadsheet(self, path: str, *, original_name: str, mime_type: str) -> ParsedFile:
        sections: list[ParsedSection] = []
        sheet_names: list[str] = []

        if pd is not None:
            workbook = pd.ExcelFile(path)
            for sheet_name in workbook.sheet_names:
                data_frame = workbook.parse(sheet_name).fillna("")
                sheet_text = data_frame.to_csv(index=False)
                cleaned = self._clean_text(sheet_text)
                if cleaned:
                    sections.append(ParsedSection(text=cleaned, sheet_name=sheet_name))
                    sheet_names.append(sheet_name)
        elif load_workbook is not None:
            workbook = load_workbook(path, data_only=True)
            for worksheet in workbook.worksheets:
                lines = []
                for row in worksheet.iter_rows(values_only=True):
                    values = [str(value).strip() for value in row if value is not None and str(value).strip()]
                    if values:
                        lines.append(" | ".join(values))
                cleaned = self._clean_text("\n".join(lines))
                if cleaned:
                    sections.append(ParsedSection(text=cleaned, sheet_name=worksheet.title))
                    sheet_names.append(worksheet.title)

        combined = "\n\n".join(
            f"{section.sheet_name}\n{section.text}" if section.sheet_name else section.text
            for section in sections
        )
        return ParsedFile(
            text=combined,
            preview_text=combined[:400],
            metadata={"parser": "pandas" if pd is not None else "openpyxl", "sheet_names": sheet_names},
            sections=sections,
        )

    async def _parse_pptx(self, path: str, *, original_name: str, mime_type: str) -> ParsedFile:
        if Presentation is None:
            raise ValueError("PPTX support is unavailable because python-pptx is not installed.")
        presentation = Presentation(path)
        sections: list[ParsedSection] = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            text_parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and str(shape.text or "").strip():
                    text_parts.append(str(shape.text).strip())
            slide_text = self._clean_text("\n".join(text_parts))
            if slide_text:
                sections.append(ParsedSection(text=slide_text, page_number=slide_index, section_title=f"Slide {slide_index}"))
        combined = "\n\n".join(section.text for section in sections)
        return ParsedFile(
            text=combined,
            preview_text=combined[:400],
            metadata={"parser": "python-pptx", "slide_count": len(sections)},
            sections=sections,
        )

    async def _parse_image(self, path: str, *, original_name: str, mime_type: str) -> ParsedFile:
        extracted_text = ocr_service.extract_text_from_path(path)
        return ParsedFile(
            text=extracted_text,
            preview_text=extracted_text[:400],
            metadata={"parser": "ocr", **ocr_service.metadata()},
            sections=[ParsedSection(text=extracted_text)] if extracted_text else [],
        )

    async def _parse_code(self, path: str, *, original_name: str, mime_type: str) -> ParsedFile:
        async with aiofiles.open(path, "r", encoding="utf-8", errors="ignore") as file_handle:
            raw_text = await file_handle.read()

        extension = Path(original_name or path).suffix.lower()
        language = CODE_LANGUAGE_BY_EXTENSION.get(extension, "code")
        lines = raw_text.splitlines()
        sections: list[ParsedSection] = []
        buffer: list[str] = []
        title = "Module"
        start_line = 1

        def flush_buffer() -> None:
            nonlocal buffer, start_line, title
            text = self._clean_text("\n".join(buffer))
            if not text:
                buffer = []
                return
            end_line = start_line + len(buffer) - 1
            sections.append(
                ParsedSection(
                    text=text,
                    section_title=f"{title} (lines {start_line}-{end_line})",
                    language=language,
                )
            )
            buffer = []

        boundary_pattern = re.compile(
            r"^\s*(?:def |class |async def |function |export function |const [A-Za-z0-9_]+\s*=\s*\(|interface |type )"
        )

        for line_number, line in enumerate(lines, start=1):
            if boundary_pattern.search(line) and buffer:
                flush_buffer()
                title = line.strip().split("{", 1)[0][:80]
                start_line = line_number
            if not buffer:
                start_line = line_number
            buffer.append(line)
            if len(buffer) >= 80:
                flush_buffer()
                title = "Continuation"
                start_line = line_number + 1

        if buffer:
            flush_buffer()

        cleaned = self._clean_text(raw_text)
        return ParsedFile(
            text=cleaned,
            preview_text=cleaned[:400],
            metadata={"parser": "code", "language": language},
            sections=sections or [ParsedSection(text=cleaned, language=language, section_title="Module")],
        )

    def _clean_text(self, text: str) -> str:
        cleaned = str(text or "").replace("\x00", " ")
        cleaned = cleaned.replace("\r\n", "\n").replace("\r", "\n")
        cleaned = re.sub(r"\n{3,}", "\n\n", cleaned)
        cleaned = re.sub(r"[ \t]{2,}", " ", cleaned)
        return cleaned.strip()


file_parser_service = FileParserService()
